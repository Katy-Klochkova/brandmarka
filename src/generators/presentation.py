"""Генерация клиентской презентации в формате PPTX в стиле Anthropic."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from core.brand import load as load_brand
from core.config import BASE_DIR, OUTPUT_DIR
from core.exporter import office_to_pdf
from core.images import prepare_logo


# Anthropic-style defaults
_DEFAULT_PRIMARY = "#1F1F1E"
_DEFAULT_SECONDARY = "#5A4BFF"
_DEFAULT_ACCENT = "#FF5B24"
_DEFAULT_BACKGROUND = "#FAFAF8"
_DEFAULT_TEXT = "#1F1F1E"


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Преобразовать HEX-цвет в RGBColor для python-pptx."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _brand_colors(brand: dict) -> dict:
    """Получить цвета бренда с антропик-фолбэками."""
    colors = brand.get("colors", {})
    return {
        "primary": _hex_to_rgb(colors.get("primary", _DEFAULT_PRIMARY)),
        "secondary": _hex_to_rgb(colors.get("secondary", _DEFAULT_SECONDARY)),
        "accent": _hex_to_rgb(colors.get("accent", _DEFAULT_ACCENT)),
        "background": _hex_to_rgb(colors.get("background", _DEFAULT_BACKGROUND)),
        "text": _hex_to_rgb(colors.get("text", _DEFAULT_TEXT)),
    }


def _set_background(slide, brand_colors: dict) -> None:
    """Установить фон слайда."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = brand_colors["background"]


def _add_logo(slide, brand: dict, brand_colors: dict) -> None:
    """Добавить логотип в правый верхний угол слайда, если он есть."""
    logo_url = brand.get("logo_url", "")
    logo_path = prepare_logo(logo_url, BASE_DIR, width=300)
    if not logo_path:
        return
    try:
        # Макс. высота 0.55 дюйма (~50 px), сохраняем пропорции
        logo = slide.shapes.add_picture(str(logo_path), Inches(9.0), Inches(0.25))
        max_height = Inches(0.55)
        aspect = logo.width / logo.height
        logo.height = max_height
        logo.width = int(max_height * aspect)
        # Сдвинуть вправо по ширине
        logo.left = Inches(10) - logo.width - Inches(0.4)
    except Exception:
        pass


def _add_accent_bar(slide, brand_colors: dict, top: bool = True) -> None:
    """Добавить тонкую цветную полосу в стиле Anthropic."""
    if top:
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    else:
        shape = slide.shapes.add_shape(1, Inches(0), Inches(5.545), Inches(10), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = brand_colors["secondary"]
    shape.line.fill.background()


def _add_title_slide(prs: Presentation, brand: dict, client: str, brand_colors: dict) -> None:
    """Титульный слайд."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_background(slide, brand_colors)

    _add_accent_bar(slide, brand_colors, top=True)
    _add_accent_bar(slide, brand_colors, top=False)
    _add_logo(slide, brand, brand_colors)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(8.6), Inches(1.2))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Презентация для {client}"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = brand_colors["primary"]

    subtitle = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(8.6), Inches(0.8))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = brand.get("company", "")
    p.font.size = Pt(26)
    p.font.color.rgb = brand_colors["text"]

    if brand.get("slogan"):
        slogan = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(8.6), Inches(0.6))
        tf = slogan.text_frame
        p = tf.paragraphs[0]
        p.text = brand["slogan"]
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = brand_colors["secondary"]


def _add_bullet_slide(prs: Presentation, brand: dict, brand_colors: dict, title_text: str, bullets: list[str]) -> None:
    """Слайд с заголовком и списком."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_background(slide, brand_colors)

    _add_accent_bar(slide, brand_colors, top=True)
    _add_logo(slide, brand, brand_colors)

    # Декор слева
    decor = slide.shapes.add_shape(1, Inches(0.55), Inches(0.85), Inches(0.08), Inches(0.5))
    decor.fill.solid()
    decor.fill.fore_color.rgb = brand_colors["accent"]
    decor.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(8.5), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = brand_colors["primary"]

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(8.2), Inches(3.6))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.size = Pt(23)
        p.font.color.rgb = brand_colors["text"]
        p.space_after = Pt(16)


def _add_contacts_slide(prs: Presentation, brand: dict, brand_colors: dict) -> None:
    """Финальный слайд с контактами."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_background(slide, brand_colors)

    # Фиолетовый блок слева
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.2), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = brand_colors["secondary"]
    bg.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(3.4), Inches(0.9))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Свяжитесь с нами"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    contacts = brand.get("contacts", {})
    lines = [
        brand.get("company", ""),
        contacts.get("phone", ""),
        contacts.get("email", ""),
        contacts.get("site", ""),
        contacts.get("address", ""),
    ]

    body = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(3.4), Inches(3.0))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate([l for l in lines if l]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(10)

    # Логотип справа, если есть
    _add_logo(slide, brand, brand_colors)


def generate(input_data: dict, output_dir: Path = OUTPUT_DIR, output_format: str = "pptx") -> Path:
    """Сгенерировать презентацию и вернуть путь к файлу."""
    brand = load_brand()
    output_dir.mkdir(parents=True, exist_ok=True)

    client_name = input_data.get("client_name", "Клиент")
    project_goal = input_data.get("project_goal", "")
    key_points = input_data.get("key_points", [])

    brand_colors = _brand_colors(brand)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    _add_title_slide(prs, brand, client_name, brand_colors)

    if project_goal:
        _add_bullet_slide(prs, brand, brand_colors, "Цель проекта", [project_goal])

    if key_points:
        _add_bullet_slide(prs, brand, brand_colors, "Почему мы", key_points)

    _add_bullet_slide(prs, brand, brand_colors, "Следующие шаги", [
        "Согласование технического задания",
        "Подготовка коммерческого предложения",
        "Демонстрация решения",
    ])

    _add_contacts_slide(prs, brand, brand_colors)

    safe_name = client_name.replace(" ", "_").replace(".", "")
    pptx_path = output_dir / f"presentation_{safe_name}.pptx"
    prs.save(pptx_path)

    if output_format.lower() == "pdf":
        return office_to_pdf(pptx_path, output_dir)
    return pptx_path


if __name__ == "__main__":
    sample = {
        "client_name": "ООО РоторПром",
        "project_goal": "внедрение системы онлайн-вибродиагностики на мельничном оборудовании",
        "key_points": [
            "15 лет опыта в диагностике вращающегося оборудования",
            "более 200 успешных внедрений",
            "собственное производство датчиков",
            "техническая поддержка 24/7",
        ],
        "pages": 8,
    }
    path = generate(sample)
    print(f"Презентация сохранена: {path}")
